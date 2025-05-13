import datetime
import json
import os
import re
import subprocess
from pathlib import Path

import markdown
import pypandoc
import yaml
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from common.logconfig import LOGGER


# Fonction récursive pour ajouter des diapositives pour chaque sous-section
def add_slides_recursively(content, presentation, level=1):
    # Diviser le contenu en sous-sections basées sur les titres de niveau actuel
    sections = content.split(f"<h{level}>")

    for i, section in enumerate(sections):
        if i == 0:
            continue  # Ignorer la première section vide

        # Créer une nouvelle diapositive
        slide_layout = presentation.slide_layouts[1]  # Layout Titre et Contenu
        slide = presentation.slides.add_slide(slide_layout)

        # Définir le titre de la diapositive
        title_end_index = section.find(f"</h{level}>")
        title = section[:title_end_index].strip()
        slide.shapes.title.text = title

        # Définir le contenu de la diapositive
        content = section[title_end_index + len(f"</h{level}>") :].strip()
        text_frame = slide.shapes.placeholders[1].text_frame
        text_frame.text = content

        # Appeler récursivement pour les sous-sections de niveau inférieur
        if level < 6:
            add_slides_recursively(content, presentation, level + 1)

        return presentation


def beamer_to_pdf(tex_file, output_file, engine="pypandoc"):

    if engine == "pypandoc":
        try:
            pypandoc.convert_file(
                tex_file,
                "beamer",
                outputfile=output_file,
                extra_args=["--pdf-engine=pdflatex"],
            )
            LOGGER.info(f"Conversion successful! Output file: {output_file}")
        except Exception as e:
            LOGGER.error(f"Conversion failed with error: {e}")
    elif engine == "pdflatex":
        subprocess.run(["pdflatex", tex_file])
    else:
        raise ValueError("not implemented")
    return output_file


def dict_to_yaml(data, filename):
    """
    Converts a dictionary to a YAML file with UTF-8 encoding.

    :param data: Dictionary to convert
    :param filename: Name of the file to write the YAML content
    """
    with open(filename, "w", encoding="utf-8") as file:
        yaml.dump(
            data, file, default_flow_style=False, sort_keys=False, allow_unicode=True
        )
    LOGGER.debug(f"Dictionnaire exported in {filename}")
    return


def dict_to_json(data, filename):
    with open(filename, "w", encoding="utf-8") as f:
        f.write(json.dumps(data, ensure_ascii=False))
    LOGGER.debug(f"Dictionnaire exported in {filename}")
    return


def clean_and_export_file(input_file_path, en_tete="```yaml", fin="```"):

    output_file_path = input_file_path.with_suffix(".yaml")

    with open(input_file_path, "r", encoding="utf-8") as file:
        lines = file.readlines()

    if lines and lines[0].strip() == en_tete:
        lines.pop(0)
    if lines and lines[-1].strip() == fin:
        lines.pop()

    # output_file_path = input_file_path.rsplit(".", 1)[0] + ".yaml"

    with open(output_file_path, "w", encoding="utf-8") as file:
        file.writelines(lines)

    LOGGER.info(f"Fichier nettoyé et exporté vers {output_file_path}")

    return output_file_path


def markdown_to_pptx(input_file, output_file):

    # Lire le contenu du fichier Markdown d'entrée
    with open(input_file, "r", encoding="utf-8") as file:
        markdown_content = file.read()

    # Convertir le contenu Markdown en HTML
    html_content = markdown.markdown(markdown_content)

    # Créer une présentation PowerPoint
    presentation = Presentation()

    # Ajouter des diapositives pour chaque section de niveau 1
    presentation = add_slides_recursively(html_content, presentation)

    # Enregistrer la présentation
    presentation.save(output_file)

    LOGGER.debug(f"Converti {input_file} en {output_file}")


def markdown_to_docx(markdown_file, word_file):
    try:
        # Convertir le fichier Markdown en LaTeX
        output = pypandoc.convert_file(markdown_file, "docx", outputfile=word_file)
        LOGGER.info(f"Conversion réussie : {word_file}")
    except Exception as e:
        LOGGER.error(f"Erreur lors de la conversion : {e}")


def markdown_to_pdf(markdown_file, pdf_file):
    try:
        # Convertir le fichier Markdown en LaTeX
        output = pypandoc.convert_file(markdown_file, "pdf", outputfile=pdf_file)
        LOGGER.info(f"Conversion réussie : {pdf_file}")
    except Exception as e:
        LOGGER.error(f"Erreur lors de la conversion : {e}")


def markdown_to_beamer(
    input_file: Path,
    output_pdf: Path = "slides.pdf",
    output_tex: Path = "slides.tex",
    title: str = "Summary",
):

    if not os.path.exists(input_file):
        LOGGER.error(f"Error: {input_file} not found.")
        return
    add_header_to_beamer_markdown(input_file, input_file, title)
    try:
        # Convert to PDF
        pypandoc.convert_file(
            input_file,
            "beamer",
            outputfile=output_pdf,
            extra_args=["--pdf-engine=xelatex"],
        )
        LOGGER.info(f"PDF conversion complete: {output_pdf}")

        # Convert to TEX
        pypandoc.convert_file(
            input_file, "beamer", outputfile=output_tex, extra_args=["--standalone"]
        )
        LOGGER.info(f"TEX conversion complete: {output_tex}")
    except Exception as e:
        LOGGER.error(f"Error during conversion: {e}")


def format_summary_as_markdown(summary):
    # Initialiser le contenu Markdown
    markdown_content = "# Résumé\n\n"

    # Ajouter le résumé au contenu Markdown
    markdown_content += summary + "\n"

    return markdown_content


def json_to_pdf(file_path, output_pdf):

    try:
        with open(file_path, "r") as file:
            data = json.load(file)
    except FileNotFoundError:
        LOGGER.error(f"Error: The file {file_path} does not exist.")
    except json.JSONDecodeError:
        LOGGER.error(f"Error: The file {file_path} is not a valid JSON file.")
    except Exception as e:
        LOGGER.error(f"An unexpected error occurred: {e}")

    # Create a canvas object
    c = canvas.Canvas(output_pdf, pagesize=letter)
    width, height = letter

    # Set initial position
    y_position = height - 40

    # Function to draw text
    def draw_text(text, x, y):
        c.drawString(x, y, text)

    # Draw JSON data
    for key, value in data.items():
        text = f"{key}: {value}"
        draw_text(text, 40, y_position)
        y_position -= 20
        if y_position < 40:
            c.showPage()
            y_position = height - 40

    # Save the PDF file
    c.save()


def json_to_yaml(json_file_path, yaml_file_path):
    # Read the JSON file
    with open(json_file_path, "r", encoding="utf-8") as json_file:
        data = json.load(json_file)

    # Convert the JSON data to YAML
    yaml_data = yaml.dump(data, allow_unicode=True, default_flow_style=False)
    # LOGGER.info(yaml_data)
    # Write the YAML data to a file
    with open(yaml_file_path, "w", encoding="utf-8") as yaml_file:
        yaml_file.write(yaml_data)

    # LOGGER.info(json_file_path, "converted to", yaml_file_path)


def yaml_to_dict(yaml_file_path):
    # Read the YAML file
    with open(yaml_file_path, "r", encoding="utf-8") as yaml_file:
        data = yaml.safe_load(yaml_file)
    # test_pappers_data_compliance(data)
    LOGGER.debug(str(yaml_file_path) + " loaded")
    return data


def add_header_to_beamer_markdown(
    markdown_file_path, beamer_mdpath, title="Résumé de l'entreprise"
):

    # je supprime les premiers headers
    remove_unique_level_headers_and_code_blocks(markdown_file_path)

    header = (
        f'---\ntitle: "'
        + title
        + '"'
        + '\nauthor: "Comptoirs et Commerces"\n'
        + 'date: "'
        + datetime.datetime.today().strftime("%Y-%m-%d")
        + '"\n---'
    )

    with open(markdown_file_path, "r", encoding="utf-8") as file:
        content = file.read()

    # Use regex to find the first header line
    pattern = r"^# .*$"
    match = re.search(pattern, content, re.MULTILINE)

    if match:
        # Remove the first header line
        content = re.sub(pattern, "", content, count=1, flags=re.MULTILINE)
        # Remove any leading newlines that might result from the removal
        # content = re.sub(r"^\n", "", content, flags=re.MULTILINE)

    # # Prepend the header

    new_content = header + "\n" + content

    # add a conclusion header before the last line
    # splitted_c = new_content.split("\n")
    # splitted_c.insert(-1, "\n## Conclusion\n")
    # new_content = "\n".join(splitted_c)

    # Use regex to replace '\n#' with '\n\n#'
    # pattern = r"\n#"
    # replacement = r"\n\n#"
    # new_content = re.sub(pattern, replacement, new_content)

    with open(beamer_mdpath, "w", encoding="utf-8") as file:
        file.write(new_content)

    LOGGER.info(f"Header added to {markdown_file_path}")
    LOGGER.info(f"Exported to {beamer_mdpath}")

    return


def add_title_to_markdown(
    markdown_file_path, title="Résumé de l'entreprise", header=False
):

    # je supprime les premiers headers
    remove_unique_level_headers_and_code_blocks(markdown_file_path)

    if header:
        header = (
            f'---\ntitle: "'
            + title
            + '"'
            + '\nauthor: "Comptoirs et Commerces"\n'
            + 'date: "'
            + datetime.datetime.today().strftime("%Y-%m-%d")
            + '"\n---'
        )
    else:
        header = f"# {title}"

    with open(markdown_file_path, "r", encoding="utf-8") as file:
        content = file.read()

    # Use regex to find the first header line
    pattern = r"^# .*$"
    match = re.search(pattern, content, re.MULTILINE)

    if match:
        # Remove the first header line
        content = re.sub(pattern, "", content, count=1, flags=re.MULTILINE)
        # Remove any leading newlines that might result from the removal
        # content = re.sub(r"^\n", "", content, flags=re.MULTILINE)

    # # Prepend the header

    new_content = header + "\n" + content

    with open(markdown_file_path, "w", encoding="utf-8") as file:
        file.write(new_content)

    LOGGER.info(f"Header added to {markdown_file_path}")
    return


def remove_unique_level_headers_and_code_blocks(file_path):
    with open(file_path, "r", encoding="utf-8") as file:
        lines = file.readlines()

    # Dictionaries to count occurr,ences of each header level
    header_count = {"#": 0, "##": 0, "###": 0, "####": 0}

    # Identify headers and count their occurrences
    for line in lines:
        if line.strip().startswith("# ") or line.strip().startswith("## "):
            level = line.strip().split(" ")[0]
            header_count[level] += 1

    # Filter out unique level headers and code blocks
    non_unique_lines = [
        line
        for line in lines
        if not (
            line.strip().startswith("#")
            and header_count[line.strip().split(" ")[0]] == 1
        )
        and not line.strip().startswith("```")
    ]

    # Write the result back to the file or a new file
    with open(file_path, "w", encoding="utf-8") as file:
        file.writelines(non_unique_lines)


def pt_to_cm(pt):
    return (pt * 1.3333343412075) * 0.0264583333


def test_pappers_data_compliance(data):
    """different test pour se convaincre de la tete de la data

    Args:
        data (_type_): _description_
    """
    LOGGER.debug(f"Societe {data['denomination']}")
    LOGGER.debug(f"siren {data['siren']}")
    if data["beneficiaires_effectifs"] != []:
        LOGGER.debug("il existe des beneficiaires effectifs")
        for beneff in data["beneficiaires_effectifs"]:
            noms = beneff["prenom_usuel"] + " " + beneff["nom"]
            LOGGER.debug(noms)
            for key in data.keys():
                if "part" in key:
                    LOGGER.debug(key, data[key])
    else:
        LOGGER.debug("pas de beneficiaires effectifs")
    return


def docx_to_pdf(docxfile_path, pdffile_path=None):
    if pdffile_path is None:
        pdffile_path = docxfile_path.with_suffix(".pdf")

    LOGGER.debug("conversion docx to pdf")
    LOGGER.debug(docxfile_path)
    LOGGER.debug(pdffile_path)

    pypandoc.convert_file(docxfile_path, "pdf", outputfile=pdffile_path)

    return pdffile_path


def main():
    # Example usage
    data = {
        "GALLA_siren": "310130323",
        "BISTROT_VALOIS_siren": "481383537",
        "LE_JARDIN_DE_ROME_siret": "57205690100018",
    }

    dict_to_yaml(data, "output.yaml")


if __name__ == "__main__":
    # main()
    # Exemple d'utilisation
    clean_and_export_file(
        Path(
            r"C:\Users\lvolat\Documents\cocoAI\work\output\bail_Annexe_6_b_-_Bail_du_29_mars_2016\extraction_results.md"
        )
    )
